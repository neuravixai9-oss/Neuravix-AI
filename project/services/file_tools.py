"""
Сборка готовых файлов для скачивания из содержимого, сгенерированного
нейросетью через функцию create_file. Каждая функция должна быть
устойчивой к ошибкам — при сбое сложного формата (docx/xlsx/pptx)
вызывающий код откатывается на обычный .txt, чтобы пользователь
в любом случае получил результат.
"""

import io


def _safe_filename(filename: str, fallback_ext: str = ".txt") -> str:
    filename = (filename or "file").strip().replace("/", "_").replace("\\", "_")
    if "." not in filename:
        filename += fallback_ext
    return filename[:120]


def build_text_file(content: str) -> bytes:
    return (content or "").encode("utf-8")


def build_docx_file(content: str) -> bytes:
    """Строки, начинающиеся с '# ', становятся заголовками; остальное — абзацы."""
    import docx
    doc = docx.Document()
    for line in (content or "").split("\n"):
        stripped = line.strip()
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:].strip(), level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:].strip(), level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:].strip(), level=1)
        else:
            doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_xlsx_file(content: str) -> bytes:
    """Строки — строки таблицы, ячейки разделены символом '|'."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    for line in (content or "").split("\n"):
        if not line.strip():
            continue
        cells = [c.strip() for c in line.split("|")]
        ws.append(cells)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_pptx_file(content: str) -> bytes:
    """Слайды разделены строкой '---', первая строка слайда — заголовок."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    title_layout = prs.slide_layouts[1]  # заголовок + текст
    raw_slides = (content or "").split("---")
    for raw in raw_slides:
        lines = [l for l in raw.strip().split("\n") if l.strip()]
        if not lines:
            continue
        title = lines[0].lstrip("#").strip()
        body = "\n".join(lines[1:])
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title[:255]
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text_frame.text = body[:2000]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_file(filename: str, file_type: str, content: str) -> tuple[bytes, str]:
    """
    Возвращает (bytes, итоговое_имя_файла).
    При ошибке генерации сложного формата — надёжно откатывается на .txt,
    чтобы пользователь не остался без результата.
    """
    file_type = (file_type or "text").lower().strip()
    try:
        if file_type == "docx":
            return build_docx_file(content), _safe_filename(filename, ".docx")
        if file_type == "xlsx":
            return build_xlsx_file(content), _safe_filename(filename, ".xlsx")
        if file_type == "pptx":
            return build_pptx_file(content), _safe_filename(filename, ".pptx")
    except Exception:
        # Откатываемся на текстовый файл, чтобы пользователь точно получил результат
        base = (filename or "file").rsplit(".", 1)[0]
        return build_text_file(content), _safe_filename(base + ".txt")

    return build_text_file(content), _safe_filename(filename, ".txt")
